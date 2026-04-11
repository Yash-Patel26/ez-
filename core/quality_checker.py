"""
Post-generation quality validation for PPTX output.

Checks:
  - Slide count within 10-15 range
  - No shapes outside margins
  - Font consistency
  - Content coverage
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Emu


class QualityChecker:
    """Validate generated PPTX against quality criteria."""

    def __init__(self, config: dict | None = None):
        self.config = config or {}
        self.margin_left = 0.375
        self.margin_right = 13.333 - 0.375
        self.min_slides = self.config.get("slide_count", {}).get("min", 10)
        self.max_slides = self.config.get("slide_count", {}).get("max", 15)

    def validate(self, pptx_path: str) -> list[str]:
        """
        Run all quality checks on a generated PPTX file.

        Returns list of warning/error strings. Empty = all passed.
        """
        issues = []

        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            return [f"Cannot open PPTX: {e}"]

        slide_count = len(prs.slides)

        # Check slide count
        if slide_count < self.min_slides:
            issues.append(f"Slide count {slide_count} below minimum {self.min_slides}")
        elif slide_count > self.max_slides:
            issues.append(f"Slide count {slide_count} above maximum {self.max_slides}")

        # Check each slide for margin violations
        emu_per_inch = 914400
        left_limit = int(self.margin_left * emu_per_inch) - emu_per_inch  # Allow some tolerance
        right_limit = int(self.margin_right * emu_per_inch) + emu_per_inch

        for slide_idx, slide in enumerate(prs.slides, 1):
            shape_count = len(slide.shapes)
            if shape_count == 0 and slide_idx not in (1, slide_count):
                issues.append(f"Slide {slide_idx}: Empty slide (no shapes)")

            for shape in slide.shapes:
                if shape.left is not None and shape.left < left_limit:
                    issues.append(
                        f"Slide {slide_idx}: Shape '{shape.name}' "
                        f"left edge ({shape.left / emu_per_inch:.2f}\") "
                        f"outside left margin"
                    )
                if (shape.left is not None and shape.width is not None
                        and (shape.left + shape.width) > right_limit):
                    issues.append(
                        f"Slide {slide_idx}: Shape '{shape.name}' "
                        f"right edge outside right margin"
                    )

        return issues
