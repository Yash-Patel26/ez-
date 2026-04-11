"""
Theme extractor for Slide Master PPTX templates.

Parses the theme XML inside a .pptx to extract:
  - Color palette (dk1, lt1, dk2, lt2, accent1-6, hlink, folHlink)
  - Font families (major/minor)
  - Slide dimensions

This ensures all generated content uses the exact design tokens
from the provided Slide Master rather than hardcoded values.
"""

from __future__ import annotations

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

from core.models import ThemeColors, ThemeFonts, ThemeConfig

# XML namespaces used in OOXML theme files
_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def _parse_color_element(el: etree._Element | None) -> str | None:
    """Extract hex color from a theme color element (srgbClr or sysClr)."""
    if el is None:
        return None

    srgb = el.find("a:srgbClr", _NS)
    if srgb is not None:
        return f"#{srgb.get('val', '000000')}"

    sys_clr = el.find("a:sysClr", _NS)
    if sys_clr is not None:
        last_clr = sys_clr.get("lastClr")
        if last_clr:
            return f"#{last_clr}"
        # Fallback to known system colors
        val = sys_clr.get("val", "")
        if val == "windowText":
            return "#000000"
        if val == "window":
            return "#FFFFFF"

    return None


def _extract_colors(theme_el: etree._Element) -> ThemeColors:
    """Extract the 12-color palette from the theme element."""
    color_scheme = theme_el.find(".//a:themeElements/a:clrScheme", _NS)
    if color_scheme is None:
        return ThemeColors()

    color_map = {}
    for name in ["dk1", "lt1", "dk2", "lt2",
                  "accent1", "accent2", "accent3", "accent4",
                  "accent5", "accent6", "hlink", "folHlink"]:
        el = color_scheme.find(f"a:{name}", _NS)
        color = _parse_color_element(el)
        if color:
            color_map[name] = color

    return ThemeColors(**color_map)


def _extract_fonts(theme_el: etree._Element) -> ThemeFonts:
    """Extract major and minor font families from the theme."""
    font_scheme = theme_el.find(".//a:themeElements/a:fontScheme", _NS)
    if font_scheme is None:
        return ThemeFonts()

    major = "Calibri"
    minor = "Calibri"

    major_el = font_scheme.find("a:majorFont/a:latin", _NS)
    if major_el is not None:
        typeface = major_el.get("typeface", "")
        if typeface and not typeface.startswith("+"):
            major = typeface

    minor_el = font_scheme.find("a:minorFont/a:latin", _NS)
    if minor_el is not None:
        typeface = minor_el.get("typeface", "")
        if typeface and not typeface.startswith("+"):
            minor = typeface

    return ThemeFonts(major=major, minor=minor)


def extract_theme(template_path: str) -> ThemeConfig:
    """
    Extract complete theme configuration from a Slide Master PPTX.

    Args:
        template_path: Path to the .pptx template file.

    Returns:
        ThemeConfig with colors, fonts, and dimensions.
    """
    prs = Presentation(template_path)

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    width_inches = slide_width / Emu(914400)  # 1 inch = 914400 EMU
    height_inches = slide_height / Emu(914400)

    # Access the theme XML from the slide master
    master = prs.slide_masters[0]
    colors = ThemeColors()
    fonts = ThemeFonts()

    try:
        # Try slide master rels first
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                theme_el = etree.fromstring(rel.target_part.blob)
                colors = _extract_colors(theme_el)
                fonts = _extract_fonts(theme_el)
                break
        else:
            # Fallback: search all package parts for theme XML
            for part in prs.part.package.iter_parts():
                if part.partname and "theme" in str(part.partname).lower():
                    theme_el = etree.fromstring(part.blob)
                    colors = _extract_colors(theme_el)
                    fonts = _extract_fonts(theme_el)
                    break
    except Exception:
        # If theme extraction fails, use defaults (Accenture palette)
        pass

    return ThemeConfig(
        colors=colors,
        fonts=fonts,
        slide_width=float(width_inches),
        slide_height=float(height_inches),
    )


def get_layout_names(template_path: str) -> list[dict]:
    """
    List all slide layouts available in the template.

    Returns list of dicts with layout index, name, and placeholder info.
    """
    prs = Presentation(template_path)
    layouts = []

    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append({
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name,
                "left": round(ph.left / 914400, 3) if ph.left else None,
                "top": round(ph.top / 914400, 3) if ph.top else None,
                "width": round(ph.width / 914400, 3) if ph.width else None,
                "height": round(ph.height / 914400, 3) if ph.height else None,
            })

        layouts.append({
            "index": idx,
            "name": layout.name,
            "placeholders": placeholders,
        })

    return layouts
