"""
Agent 5: Visual Generator

Creates charts, tables, and infographics using python-pptx shapes.
All visuals are programmatically generated (no external images),
using theme colors for consistency.

Responsibilities:
  - Native python-pptx charts (bar, pie, line, area)
  - Styled tables with theme-colored headers
  - KPI card shapes
  - Process flow connected shapes
  - Timeline shapes

Input:  Slide object + OptimizedSlideContent + ThemeConfig
Output: Shapes added directly to the slide
"""

from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from core.models import (
    ChartData,
    KPIItem,
    OptimizedSlideContent,
    ProcessStep,
    ShapeSpec,
    TableData,
    ThemeConfig,
    TimelineItem,
)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _boost_saturation(r: int, g: int, b: int) -> tuple[int, int, int]:
    """Boost saturation of faded/washed-out accent colors.

    Common Mistake #23: colors too faded. If an accent color has
    high luminance (>0.75) and low saturation, push saturation up
    so it reads well on white slide backgrounds.
    """
    # Relative luminance (simplified)
    lum = (0.299 * r + 0.587 * g + 0.114 * b) / 255.0
    max_c = max(r, g, b)
    min_c = min(r, g, b)
    saturation = (max_c - min_c) / max(max_c, 1)

    # Only boost if the color is both faded AND washed-out
    if lum > 0.75 and saturation < 0.35:
        # Reduce each channel toward zero proportionally to darken
        factor = 0.72
        r = int(r * factor)
        g = int(g * factor)
        b = int(b * factor)

    return r, g, b


def _resolve_color(color_ref: str, theme: ThemeConfig) -> RGBColor:
    """Resolve a color reference (theme name or hex) to RGBColor."""
    if not color_ref:
        return RGBColor(0, 0, 0)

    # References that should NEVER be boosted (text & background colors)
    _no_boost = {"dk1", "lt1", "dk2", "lt2"}

    # Check if it's a theme reference
    theme_map = {
        "dk1": theme.colors.dk1,
        "lt1": theme.colors.lt1,
        "dk2": theme.colors.dk2,
        "lt2": theme.colors.lt2,
        "accent1": theme.colors.accent1,
        "accent2": theme.colors.accent2,
        "accent3": theme.colors.accent3,
        "accent4": theme.colors.accent4,
        "accent5": theme.colors.accent5,
        "accent6": theme.colors.accent6,
        "primary": theme.colors.primary_accent(theme.primary_accent_seed),
    }

    if color_ref in theme_map:
        rgb = _hex_to_rgb(theme_map[color_ref])
        # Boost saturation for accent colors only
        if color_ref not in _no_boost:
            r, g, b = _boost_saturation(rgb[0], rgb[1], rgb[2])
            return RGBColor(r, g, b)
        return rgb

    if color_ref.startswith('#'):
        return _hex_to_rgb(color_ref)

    return RGBColor(0, 0, 0)


_C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _enable_shrink_to_fit(text_frame) -> None:
    """Enable PowerPoint's 'Shrink text on overflow' on a text frame.

    Layout engine already picks font sizes that fit the estimated text
    length, but glyph widths vary and some bullets may still overflow.
    normAutofit tells PowerPoint to auto-scale the text at render time so
    content never clips out of the box — this is the in-product 'auto-fit'
    behaviour you get by right-clicking → Size & Properties → Shrink text
    on overflow.
    """
    from lxml import etree

    bodyPr = text_frame._txBody.find(f"{{{_A_NS}}}bodyPr")
    if bodyPr is None:
        return
    # Remove any existing autofit element before adding normAutofit
    for tag in ("normAutofit", "spAutoFit", "noAutofit"):
        existing = bodyPr.find(f"{{{_A_NS}}}{tag}")
        if existing is not None:
            bodyPr.remove(existing)
    etree.SubElement(bodyPr, f"{{{_A_NS}}}normAutofit")


def _set_white_fill_xml(parent_elem) -> None:
    """Inject <c:spPr> with a solid white fill and no border into a chart XML element.

    python-pptx 1.x does not expose chart_area / plot_area as Python objects,
    so we manipulate the underlying lxml element directly.
    The resulting XML fragment:
        <c:spPr>
          <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
          <a:ln><a:noFill/></a:ln>
        </c:spPr>
    """
    from lxml import etree

    # Remove any existing spPr to avoid duplicates
    existing = parent_elem.find(f"{{{_C_NS}}}spPr")
    if existing is not None:
        parent_elem.remove(existing)

    spPr = etree.SubElement(parent_elem, f"{{{_C_NS}}}spPr")
    solidFill = etree.SubElement(spPr, f"{{{_A_NS}}}solidFill")
    srgbClr = etree.SubElement(solidFill, f"{{{_A_NS}}}srgbClr")
    srgbClr.set("val", "FFFFFF")
    ln = etree.SubElement(spPr, f"{{{_A_NS}}}ln")
    etree.SubElement(ln, f"{{{_A_NS}}}noFill")


def _disable_line_smoothing(chart) -> None:
    """Set smooth=0 on every line series to prevent bezier blurring."""
    from lxml import etree

    for ser_elem in chart.plots[0]._element.iter(f"{{{_C_NS}}}ser"):
        smooth_elem = ser_elem.find(f"{{{_C_NS}}}smooth")
        if smooth_elem is None:
            smooth_elem = etree.SubElement(ser_elem, f"{{{_C_NS}}}smooth")
        smooth_elem.set("val", "0")


class VisualGenerator:
    """Generate visual elements on slides using python-pptx."""

    def __init__(self, theme: ThemeConfig, config: dict | None = None):
        self.theme = theme
        self.config = config or {}
        self.accent_colors = theme.colors.accent_list()
        # Chart colors: prefer darker, more visible accent colors.
        # Some templates use near-white colors for accent1/4 (background accents).
        # Sort by relative luminance ascending so the darkest colors come first.
        self._chart_colors = self._build_chart_colors()

    def _build_chart_colors(self) -> list[str]:
        """Return accent colors sorted darkest-first, with near-white colors filtered out.

        Some templates use near-white hex values for accent1/accent4 as layout
        background tints.  Those colors produce invisible bars on a white chart
        background.  We sort by relative luminance (ascending) and drop any color
        whose luminance stays above 0.70 even after the saturation boost.
        """
        _FALLBACK = ["#2563EB", "#DC2626", "#16A34A", "#D97706", "#7C3AED", "#0891B2"]

        def _lum(hex_color: str) -> float:
            h = hex_color.lstrip('#')
            r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
            r2, g2, b2 = _boost_saturation(r, g, b)
            return (0.299 * r2 + 0.587 * g2 + 0.114 * b2) / 255.0

        sorted_colors = sorted(self.accent_colors, key=_lum)
        # Keep only colors that are dark enough to be visible on white
        visible = [c for c in sorted_colors if _lum(c) < 0.70]
        if not visible:
            return _FALLBACK
        # Pad with fallbacks if fewer than 3 visible accent colors
        if len(visible) < 3:
            visible = visible + [c for c in _FALLBACK if c not in visible]
        return visible

    def add_chart(
        self, slide, chart_data: ChartData, position: ShapeSpec
    ) -> None:
        """Add a native python-pptx chart to the slide."""
        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "pie": XL_CHART_TYPE.PIE,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "area": XL_CHART_TYPE.AREA,
            "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
        }

        xl_type = chart_type_map.get(chart_data.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        if not chart_data.series or not chart_data.categories:
            return  # Skip chart if no data

        data = CategoryChartData()
        data.categories = chart_data.categories

        for series in chart_data.series:
            data.add_series(series.name, series.values)

        p = position.position
        chart_shape = slide.shapes.add_chart(
            xl_type,
            Inches(p.left), Inches(p.top),
            Inches(p.width), Inches(p.height),
            data,
        )

        chart = chart_shape.chart

        # Solid white backgrounds via XML — python-pptx 1.x does not expose
        # chart_area / plot_area as objects, so we inject <c:spPr> directly.
        # This prevents slide gradient / image bleeding through the chart frame,
        # which is the primary cause of the "blurred" chart appearance.
        cs = chart._chartSpace
        _set_white_fill_xml(cs)   # chart area (outer frame)

        if chart_data.chart_type != "pie":
            plot_area = cs.find(
                f"{{{_C_NS}}}chart"
            ).find(f"{{{_C_NS}}}plotArea")
            if plot_area is not None:
                _set_white_fill_xml(plot_area)

        # Chart title if provided
        if chart_data.title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = chart_data.title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
            chart.chart_title.text_frame.paragraphs[0].font.bold = True
        else:
            chart.has_title = False

        chart.has_legend = len(chart_data.series) > 1

        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(9)

        # Style the chart with theme colors (darkest/most visible first)
        for i, series in enumerate(chart.series):
            color = self._chart_colors[i % len(self._chart_colors)]
            rgb_color = _hex_to_rgb(color)

            if chart_data.chart_type in ("line",):
                # Line charts: set line color + weight; fill doesn't apply to lines
                series.format.line.color.rgb = rgb_color
                series.format.line.width = Pt(2.25)  # crisp, visible line
            else:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = rgb_color

            # Add data labels for all chart types
            series.has_data_labels = True
            series.data_labels.font.size = Pt(9)
            series.data_labels.font.color.rgb = _hex_to_rgb(self.theme.colors.dk2)
            if chart_data.chart_type == "pie":
                series.data_labels.number_format = '0.0%'
                series.data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END
            elif chart_data.chart_type in ("bar", "stacked_bar"):
                series.data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END
            elif chart_data.chart_type in ("line", "area"):
                series.data_labels.label_position = XL_LABEL_POSITION.ABOVE

        # Disable bezier smoothing on line charts — smooth curves look blurry/imprecise
        if chart_data.chart_type == "line":
            _disable_line_smoothing(chart)

        # Style axes if present
        if chart_data.chart_type != "pie":
            try:
                cat_axis = chart.category_axis
                cat_axis.tick_labels.font.size = Pt(9)
                _axis_color = _resolve_color("lt2", self.theme)
                cat_axis.format.line.color.rgb = _axis_color
                cat_axis.format.line.width = Pt(0.75)

                val_axis = chart.value_axis
                val_axis.tick_labels.font.size = Pt(9)
                val_axis.has_major_gridlines = True
                val_axis.major_gridlines.format.line.color.rgb = _hex_to_rgb(self.theme.colors.lt2)
                val_axis.major_gridlines.format.line.width = Pt(0.5)
                val_axis.format.line.color.rgb = _axis_color
                val_axis.format.line.width = Pt(0.75)
            except Exception:
                pass

    def add_table(
        self, slide, table_data: TableData, position: ShapeSpec
    ) -> None:
        """Add a styled table to the slide."""
        n_rows = len(table_data.rows) + 1  # +1 for header
        n_cols = len(table_data.headers)

        if n_rows < 2 or n_cols < 1:
            return

        p = position.position
        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(p.left), Inches(p.top),
            Inches(p.width), Inches(p.height),
        )
        table = table_shape.table

        # Set column widths equally
        col_width = Inches(p.width / n_cols)
        for i in range(n_cols):
            table.columns[i].width = int(col_width)

        # Header row
        for i, header in enumerate(table_data.headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex_to_rgb(self.theme.colors.primary_accent(self.theme.primary_accent_seed))

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.bold = True
                paragraph.font.color.rgb = _hex_to_rgb(self.theme.colors.lt1)
                paragraph.alignment = PP_ALIGN.CENTER

        # Data rows
        for row_idx, row_data in enumerate(table_data.rows):
            for col_idx in range(n_cols):
                cell = table.cell(row_idx + 1, col_idx)
                raw_text = row_data[col_idx] if col_idx < len(row_data) else ""

                # Trend indicator: detect percentages/numbers with growth keywords
                trend_arrow = self._detect_trend_indicator(raw_text)
                cell.text = f"{trend_arrow} {raw_text}" if trend_arrow else raw_text

                # Alternating row colors — use theme lt1/lt2 instead of hardcoded grays
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _hex_to_rgb(self.theme.colors.lt2)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _hex_to_rgb(self.theme.colors.lt1)

                # Common Mistake #15: Text should be middle-aligned in tables
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.alignment = PP_ALIGN.CENTER

                    # Color-code cells with trend indicators using theme accent colors
                    if trend_arrow == "▲":
                        paragraph.font.color.rgb = _resolve_color("accent3", self.theme)
                    elif trend_arrow == "▼":
                        paragraph.font.color.rgb = _resolve_color("accent1", self.theme)
                    else:
                        paragraph.font.color.rgb = _hex_to_rgb(self.theme.colors.dk2)

                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.04)
                cell.margin_bottom = Inches(0.04)

        # Highlight row if specified
        if table_data.highlight_row is not None:
            hr = table_data.highlight_row + 1  # Account for header
            if hr < n_rows:
                for col_idx in range(n_cols):
                    cell = table.cell(hr, col_idx)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _hex_to_rgb(self.theme.colors.accent2)

    def add_kpi_cards(
        self, slide, kpis: list[KPIItem], shapes: list[ShapeSpec]
    ) -> None:
        """Add KPI metric cards to the slide."""
        # KPI shapes are added by the general shape renderer
        # This method handles any additional KPI-specific styling
        pass

    def add_shape(
        self, slide, spec: ShapeSpec, theme: ThemeConfig
    ) -> None:
        """Add a generic shape to the slide based on ShapeSpec."""
        from pptx.enum.shapes import MSO_SHAPE

        p = spec.position
        left = Inches(p.left)
        top = Inches(p.top)
        width = Inches(p.width)
        height = Inches(p.height)

        if spec.shape_type == "text_box":
            txBox = slide.shapes.add_textbox(left, top, width, height)
            self._format_textbox(txBox, spec, theme)
            _enable_shrink_to_fit(txBox.text_frame)

        elif spec.shape_type in ("rounded_rect", "rectangle"):
            shape_type = (MSO_SHAPE.ROUNDED_RECTANGLE
                          if spec.shape_type == "rounded_rect"
                          else MSO_SHAPE.RECTANGLE)
            shape = slide.shapes.add_shape(shape_type, left, top, width, height)
            self._format_shape(shape, spec, theme)

        elif spec.shape_type == "oval":
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
            self._format_shape(shape, spec, theme)

        elif spec.shape_type == "line":
            # Vertical line when width == 0, horizontal otherwise
            if spec.position.width == 0:
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top, Inches(0.03), height
                )
            else:
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.03)
                )
            if spec.font_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = _resolve_color(spec.font_color, theme)
            shape.line.fill.background()

        elif spec.shape_type == "arrow":
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _resolve_color(
                spec.font_color or "dk2", theme
            )
            shape.line.fill.background()

        elif spec.shape_type == "triangle":
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height
            )
            self._format_shape(shape, spec, theme)

    @staticmethod
    def _detect_trend_indicator(text: str) -> str | None:
        """Detect if a table cell value implies a trend and return an arrow."""
        import re
        text_lower = text.lower().strip()

        # Explicit positive indicators
        if any(w in text_lower for w in ('+', 'increase', 'growth', 'rise', 'up ', 'grew')):
            return "▲"
        # Explicit negative indicators
        if any(w in text_lower for w in ('decrease', 'decline', 'drop', 'down ', 'fell', 'loss')):
            return "▼"
        # Percentage with + prefix
        if re.match(r'^\+\d', text.strip()):
            return "▲"
        # Percentage with - prefix (negative)
        if re.match(r'^-\d', text.strip()):
            return "▼"

        return None

    def _format_textbox(self, txBox, spec: ShapeSpec, theme: ThemeConfig) -> None:
        """Apply formatting to a text box.

        Guideline 2B: Title uses MAJOR font, body uses MINOR font.
        Guideline 4A: Consistent font sizes across hierarchy.
        """
        tf = txBox.text_frame
        tf.word_wrap = True

        # Vertical anchor — default text_frame is top-aligned, so bullet rows
        # where the label is centred with a number circle would drift to the top
        # of their box. Honour the spec's vertical_alignment.
        anchor_map = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }
        if spec.vertical_alignment in anchor_map:
            tf.vertical_anchor = anchor_map[spec.vertical_alignment]

        # Common Mistake #14: "No margins within text boxes without fill color"
        # Only add internal margins if the text box has a fill color
        if spec.fill_color:
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
        else:
            tf.margin_left = Inches(0)
            tf.margin_right = Inches(0)
            tf.margin_top = Inches(0)
            tf.margin_bottom = Inches(0)

        # Determine if this is a title (uses major font) or body (minor font)
        is_title = spec.font_size and spec.font_size >= 28 and spec.position.top < 1.2
        font_name = theme.fonts.major if is_title else theme.fonts.minor

        lines = spec.text.split('\n') if spec.text else [""]
        for i, line in enumerate(lines):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()

            # Bullet formatting
            is_bullet = line.startswith('• ')
            if is_bullet:
                para.text = line[2:]
                para.level = 0
                para.space_after = Pt(8)   # More spacing between bullets (Guideline 5)
                para.space_before = Pt(2)
            else:
                para.text = line
                para.space_after = Pt(4)

            para.font.size = Pt(spec.font_size or 14)
            para.font.bold = spec.font_bold
            para.font.name = font_name

            if spec.font_color:
                para.font.color.rgb = _resolve_color(spec.font_color, theme)

            para.alignment = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
            }.get(spec.alignment, PP_ALIGN.LEFT)

    def _format_shape(self, shape, spec: ShapeSpec, theme: ThemeConfig) -> None:
        """Apply formatting to a shape.

        Guideline 4C: Clean containers, minimal borders, consistent corner radius.
        Guideline 5: Internal padding for breathing space.
        """
        # Fill
        if spec.fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _resolve_color(spec.fill_color, theme)
        else:
            shape.fill.background()

        # Border
        if spec.border_color:
            shape.line.color.rgb = _resolve_color(spec.border_color, theme)
            shape.line.width = Pt(spec.border_width or 1.0)
        else:
            shape.line.fill.background()  # No border

        # Text inside shapes
        if spec.text:
            tf = shape.text_frame
            tf.word_wrap = True
            _enable_shrink_to_fit(tf)

            # Internal padding (Guideline 5: breathing space)
            tf.margin_left = Inches(0.12)
            tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.08)
            tf.margin_bottom = Inches(0.08)

            anchor_map = {
                "top": MSO_ANCHOR.TOP,
                "middle": MSO_ANCHOR.MIDDLE,
                "bottom": MSO_ANCHOR.BOTTOM,
            }
            if spec.vertical_alignment in anchor_map:
                tf.vertical_anchor = anchor_map[spec.vertical_alignment]
            else:
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            lines = spec.text.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()

                # Handle bullets inside shapes
                is_bullet = line.startswith('• ')
                if is_bullet:
                    para.text = line[2:]
                    para.space_after = Pt(6)
                else:
                    para.text = line

                para.alignment = {
                    "left": PP_ALIGN.LEFT,
                    "center": PP_ALIGN.CENTER,
                    "right": PP_ALIGN.RIGHT,
                }.get(spec.alignment, PP_ALIGN.CENTER)

                # KPI-style: first line large value, second line small label
                if i == 0 and spec.font_size and spec.font_size >= 28:
                    para.font.size = Pt(spec.font_size)
                    para.font.bold = True
                    if spec.font_color:
                        para.font.color.rgb = _resolve_color(spec.font_color, theme)
                elif spec.font_bold and spec.font_size:
                    para.font.size = Pt(spec.font_size)
                    para.font.bold = spec.font_bold
                    if spec.font_color:
                        para.font.color.rgb = _resolve_color(spec.font_color, theme)
                else:
                    para.font.size = Pt(spec.font_size or 12)
                    para.font.color.rgb = _resolve_color(
                        spec.font_color or "dk2", theme
                    )

                para.font.name = theme.fonts.minor
