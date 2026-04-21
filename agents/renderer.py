"""
Agent 6: PPTX Renderer

Assembles the final .pptx file by combining layouts, optimized content,
and visual elements with the Slide Master template.

Responsibilities:
  - Load Slide Master template as base presentation
  - Add slides using correct layouts
  - Place text, shapes, charts, and tables via VisualGenerator
  - Apply theme fonts and colors consistently
  - Handle cover and thank-you slides via template placeholders

Input:  List[SlideLayout] + List[OptimizedSlideContent] + template path
Output: Saved .pptx file
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Emu, Inches, Pt
from pptx.enum.text import PP_ALIGN

from agents.visual_generator import VisualGenerator, _resolve_color
from core.models import (
    OptimizedSlideContent,
    ShapeSpec,
    SlideLayout,
    ThemeConfig,
)


def _A(tag: str) -> str:
    return f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}"


def _placeholder_over_dark(slide, ph) -> bool:
    """Return True if a dark-filled shape in the layout covers the placeholder."""
    left = ph.left or 0
    top = ph.top or 0
    right = left + (ph.width or 0)
    bottom = top + (ph.height or 0)
    for sh in slide.slide_layout.shapes:
        try:
            sh_left = sh.left or 0
            sh_top = sh.top or 0
            sh_right = sh_left + (sh.width or 0)
            sh_bottom = sh_top + (sh.height or 0)
            if sh_left > left or sh_top > top or sh_right < right or sh_bottom < bottom:
                continue
            rgb = sh.fill.fore_color.rgb
        except Exception:
            continue
        r, g, b = rgb[0], rgb[1], rgb[2]
        if (r * 0.299 + g * 0.587 + b * 0.114) < 128:
            return True
    return False


def _read_layout_title_style(slide, idx: int) -> tuple[int | None, RGBColor | None]:
    """Extract the layout placeholder's intended font size (pt) and RGB color.

    Looks in rPr and lstStyle defRPr. Returns (None, None) if not defined —
    caller then falls back to config defaults.
    """
    layout = slide.slide_layout
    target = None
    for lp in layout.placeholders:
        try:
            if lp.placeholder_format.idx == idx:
                target = lp
                break
        except Exception:
            continue
    if target is None:
        return None, None

    size_pt: int | None = None
    color: RGBColor | None = None

    for el in target._element.iter():
        if el.tag in (_A("rPr"), _A("defRPr")):
            sz = el.get("sz")
            if sz and size_pt is None:
                size_pt = int(sz) // 100
            for child in el:
                if child.tag == _A("solidFill") and color is None:
                    rgb_el = child.find(_A("srgbClr"))
                    if rgb_el is not None and rgb_el.get("val"):
                        color = RGBColor.from_string(rgb_el.get("val"))
    return size_pt, color


class PPTXRenderer:
    """Render the final PPTX presentation."""

    def __init__(
        self,
        template_path: str,
        theme: ThemeConfig,
        config: dict | None = None,
    ):
        self.template_path = template_path
        self.theme = theme
        self.config = config or {}
        self.visual_gen = VisualGenerator(theme, config)

    def render(
        self,
        layouts: list[SlideLayout],
        slides_content: list[OptimizedSlideContent],
        output_path: str,
    ) -> None:
        """
        Render all slides and save the presentation.

        Args:
            layouts: Computed layout specifications.
            slides_content: Optimized content for each slide.
            output_path: Where to save the .pptx file.
        """
        prs = Presentation(self.template_path)

        # Remove existing slides from template (keep only layouts/master)
        self._remove_existing_slides(prs)

        # Strip prompt/citation text from layout placeholders so empty slide-
        # level placeholders don't render the layout's baked-in text
        # (e.g. 'Source', 'Click to edit Master title style').
        self._scrub_layout_prompt_text(prs)

        # Build a layout name → layout object mapping
        layout_map = self._build_layout_map(prs)

        # Add each slide
        for layout_spec, content in zip(layouts, slides_content):
            pptx_layout = self._find_layout(layout_map, layout_spec.layout_name)
            slide = prs.slides.add_slide(pptx_layout)

            # Render based on slide type
            if content.slide_type == "cover":
                self._render_cover(slide, content, layout_spec)
            elif content.slide_type == "thank_you":
                self._render_thank_you(slide, content)
            elif content.slide_type == "divider":
                self._render_divider(slide, layout_spec, content)
            else:
                self._render_content_slide(slide, layout_spec, content)

        # Ensure output directory exists
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

    def _remove_existing_slides(self, prs: Presentation) -> None:
        """Remove all existing slides from the presentation."""
        from lxml import etree
        ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

        while len(prs.slides) > 0:
            sldId = prs.slides._sldIdLst[0]
            rId = sldId.get(f'{ns}id') or sldId.get('r:id')
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except KeyError:
                    pass
            prs.slides._sldIdLst.remove(sldId)

    def _build_layout_map(self, prs: Presentation) -> dict[str, object]:
        """Build a mapping from layout name to layout object.

        Templates often have duplicate layout names (e.g. three copies of
        '1_E_Title, Subtitle and Body', where the last one is a 'Thank you!'
        variant with that text baked into its shapes).  Overwriting the map
        with the last duplicate would inject 'Thank you!' onto every content
        slide.  Fix: store only the FIRST occurrence under the plain name, and
        separately capture any layout whose shapes contain 'thank you' text so
        the closing slide can find it explicitly.
        """
        layout_map = {}
        for layout in prs.slide_layouts:
            idx_key = f"idx_{prs.slide_layouts.index(layout)}"
            layout_map[idx_key] = layout

            name_lower = layout.name.lower()

            # Detect layouts that have 'thank you' text baked into their shapes
            has_thankyou_text = any(
                "thank you" in sh.text_frame.text.lower()
                for sh in layout.shapes
                if sh.has_text_frame
            )
            if has_thankyou_text:
                # Store under a special key for the closing slide, but do NOT
                # overwrite the clean content-layout entry under the plain name
                layout_map.setdefault("_thankyou_layout", layout)
                continue

            # Only store the first occurrence under the plain name
            layout_map.setdefault(name_lower, layout)

        return layout_map

    def _find_layout(self, layout_map: dict, layout_name: str) -> object:
        """Find the best matching layout by name."""
        name_lower = layout_name.lower()

        # Direct match
        if name_lower in layout_map:
            return layout_map[name_lower]

        # Try common names based on type
        type_names = {
            "cover": ["1_cover", "2_cover", "cover", "0_title company",
                       "title company", "title slide"],
            "divider": ["divider", "c_section blue"],
            "blank": ["blank", "1_e_title, subtitle and body"],
            "title_only": ["title only", "1_e_title, subtitle and body"],
            # Prefer the layout that already has 'thank you' text baked in
            "thank_you": ["_thankyou_layout", "thank you", "1_thank you",
                          "0_title company", "title company"],
        }

        if name_lower in type_names:
            for candidate in type_names[name_lower]:
                if candidate in layout_map:
                    return layout_map[candidate]

        # Fuzzy match
        for key in layout_map:
            if name_lower in key or key in name_lower:
                return layout_map[key]

        # Fallback: try title-related, then first layout with placeholders, then first
        for fallback in ["blank", "title only", "1_e_title, subtitle and body"]:
            if fallback in layout_map:
                return layout_map[fallback]

        # Last resort: first layout
        return list(layout_map.values())[0]

    def _render_cover(
        self, slide, content: OptimizedSlideContent, layout
    ) -> None:
        """Render the cover slide using template placeholders + decorative shapes."""
        placeholders = list(slide.placeholders)

        # ph.text wipes layout run-level rPr. Preserve the layout's intended
        # size/color for the cover so it stays visible on the template's bg.
        title_fallback = int(self.config.get("typography", {}).get("cover_title_size", 32))
        subtitle_fallback = int(self.config.get("typography", {}).get("cover_subtitle_size", 16))

        def _apply(ph, text, is_title, fallback_size):
            layout_idx = ph.placeholder_format.idx
            layout_size, layout_color = _read_layout_title_style(slide, layout_idx)
            # Guarantee contrast — if layout says white text but no dark bg shape
            # sits behind the placeholder, swap to dk1 so it's not invisible.
            if (
                layout_color
                and layout_color == RGBColor(0xFF, 0xFF, 0xFF)
                and not _placeholder_over_dark(slide, ph)
            ):
                layout_color = _resolve_color("dk1", self.theme)
            final_color = layout_color or _resolve_color(
                "dk1" if is_title else "dk2", self.theme
            )
            final_size = layout_size or fallback_size

            ph.text = text
            tf = ph.text_frame
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.word_wrap = True
            for para in tf.paragraphs:
                para.font.name = (
                    self.theme.fonts.major if is_title else self.theme.fonts.minor
                )
                para.font.size = Pt(final_size)
                para.font.bold = is_title
                para.font.color.rgb = final_color

        if len(placeholders) >= 2:
            _apply(placeholders[0], content.title or "", True, title_fallback)
            if content.subtitle:
                _apply(placeholders[1], content.subtitle, False, subtitle_fallback)
        elif len(placeholders) == 1:
            _apply(placeholders[0], content.title or "", True, title_fallback)
        else:
            # No placeholders — add text boxes
            txBox = slide.shapes.add_textbox(
                Inches(0.375), Inches(3.1), Inches(9.3), Inches(0.7)
            )
            tf = txBox.text_frame
            tf.paragraphs[0].text = content.title or ""
            tf.paragraphs[0].font.size = Pt(36)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.name = self.theme.fonts.major
            tf.paragraphs[0].font.color.rgb = _resolve_color("dk1", self.theme)

            if content.subtitle:
                txBox2 = slide.shapes.add_textbox(
                    Inches(0.375), Inches(4.3), Inches(6.6), Inches(0.8)
                )
                tf2 = txBox2.text_frame
                tf2.paragraphs[0].text = content.subtitle
                tf2.paragraphs[0].font.size = Pt(16)
                tf2.paragraphs[0].font.name = self.theme.fonts.minor
                tf2.paragraphs[0].font.color.rgb = _resolve_color("dk2", self.theme)

        # Add the layout engine's decorative shapes on top of placeholders.
        # Skip pure text-box shapes (title/subtitle already handled above)
        # to avoid duplicates; only add non-text decorative elements.
        for spec in getattr(layout, 'shapes', []):
            if spec.shape_type == "text_box":
                continue  # title/subtitle from placeholders already
            self.visual_gen.add_shape(slide, spec, self.theme)

    def _render_thank_you(self, slide, content: OptimizedSlideContent) -> None:
        """Render the thank you/closing slide.

        If the template layout already has 'Thank you' text baked into its
        shapes, we leave it alone — adding text would create a duplicate.
        Only populate placeholders or add a fallback text box when the layout
        has no such pre-existing text.
        """
        # Check if the layout already has "thank you" text baked in
        layout_has_thankyou = any(
            "thank you" in sh.text_frame.text.lower()
            for sh in slide.slide_layout.shapes
            if sh.has_text_frame
        )
        # Some templates (e.g. AI Bubble) render 'Thank you' as a vector
        # shape/picture rather than live text. Trust the layout name in that
        # case so we don't stamp a duplicate textbox on top of the artwork.
        if not layout_has_thankyou:
            layout_name = (slide.slide_layout.name or "").lower()
            if "thank you" in layout_name or "thankyou" in layout_name:
                layout_has_thankyou = True
        if layout_has_thankyou:
            # Template provides the text — nothing to add
            return

        placeholders = list(slide.placeholders)

        if placeholders:
            # Use template placeholders
            placeholders[0].text = "Thank You"
            for ph in placeholders:
                ph.text_frame.margin_left = 0
                ph.text_frame.margin_right = 0
                for para in ph.text_frame.paragraphs:
                    para.font.name = self.theme.fonts.major
        else:
            # Fallback: add "Thank You" text box so slide is never blank
            txBox = slide.shapes.add_textbox(
                Inches(2.0), Inches(2.5), Inches(9.3), Inches(1.0)
            )
            tf = txBox.text_frame
            tf.paragraphs[0].text = "Thank You"
            tf.paragraphs[0].font.size = Pt(44)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.name = self.theme.fonts.major
            tf.paragraphs[0].font.color.rgb = _resolve_color("dk1", self.theme)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _render_divider(
        self, slide, layout: SlideLayout, content: OptimizedSlideContent
    ) -> None:
        """Render a section-divider slide.

        The C_Section blue template layout provides the dark background.
        We populate the title placeholder so PowerPoint's outline/sorter views
        show the section name correctly, then add decorative shapes on top.
        """
        filled = self._populate_placeholders(slide, content.title)
        self._strip_unused_placeholders(slide, keep_idxs=filled)
        for spec in layout.shapes:
            self.visual_gen.add_shape(slide, spec, self.theme)

    def _scrub_layout_prompt_text(self, prs) -> None:
        """Empty out literal prompt/citation text on layout placeholders.

        Some templates bake prompt text directly into layout placeholders
        (e.g. Accenture's 'Title only' layout has 'Click to edit Master title
        style' on the title PH and 'Source' on a body PH). When a slide using
        that layout leaves the PH empty, PowerPoint renders the *layout*'s
        text as a fallback — bleeding the prompt through on every slide.
        Clearing the layout's placeholder text fixes the bleed-through for
        every slide in one pass without changing the master.
        """
        for layout in prs.slide_layouts:
            for ph in layout.placeholders:
                if not ph.has_text_frame:
                    continue
                try:
                    ptype = ph.placeholder_format.type
                except Exception:
                    ptype = None
                # Leave slide-number / footer / date untouched — these
                # legitimately carry auto-filled values like '<#>' or dates.
                if ptype in (
                    PP_PLACEHOLDER.SLIDE_NUMBER,
                    PP_PLACEHOLDER.FOOTER,
                    PP_PLACEHOLDER.DATE,
                ):
                    continue
                # Clear text by replacing the <a:txBody> paragraphs with a
                # single empty paragraph. ph.text = "" triggers the python-
                # pptx autoclean path that preserves defRPr style but removes
                # run-level text, which is what we want.
                ph.text_frame.clear()
                try:
                    ph.text_frame.paragraphs[0].text = ""
                except Exception:
                    pass

    def _strip_unused_placeholders(self, slide, keep_idxs: set[int]) -> None:
        """Remove placeholders whose inherited layout text would show through.

        Layouts like Accenture's 'Title only' carry a 'Source' body placeholder
        (idx 11) and 'Click to edit Master title style' prompts. When the slide's
        placeholder is empty, PowerPoint renders the layout's text — overlapping
        the content shapes we add on top. Remove every placeholder we're not
        populating, except slide-number / footer / date (harmless chrome).
        """
        keep_types = {
            PP_PLACEHOLDER.SLIDE_NUMBER,
            PP_PLACEHOLDER.FOOTER,
            PP_PLACEHOLDER.DATE,
        }
        to_remove = []
        for ph in slide.placeholders:
            try:
                idx = ph.placeholder_format.idx
                ptype = ph.placeholder_format.type
            except Exception:
                continue
            if idx in keep_idxs or ptype in keep_types:
                continue
            to_remove.append(ph)

        for ph in to_remove:
            sp = ph._element
            parent = sp.getparent()
            if parent is not None:
                parent.remove(sp)

    def _populate_placeholders(self, slide, title: str, subtitle: str = "") -> set[int]:
        """Set title and subtitle template placeholders.

        ph.text wipes the layout's run-level rPr (size, color), so the text
        falls back to master style defaults — typically 44pt title, which
        overflows a 0.57in title band. Read the layout's intended size/color
        first, then re-apply them explicitly. Size is also clamped to fit the
        placeholder's height so long titles never overflow or wrap into the
        header divider line.

        Returns the set of placeholder idxs that got populated — the caller
        uses this to know which placeholders to keep vs strip.
        """
        title_fallback = int(self.config.get("typography", {}).get("title_size", 28))
        subtitle_fallback = int(self.config.get("typography", {}).get("subtitle_size", 13))

        filled: set[int] = set()

        for ph in slide.placeholders:
            try:
                idx = ph.placeholder_format.idx
                ptype = ph.placeholder_format.type
            except Exception:
                continue

            is_subtitle_ph = idx == 1 or ptype == PP_PLACEHOLDER.SUBTITLE

            if idx == 0 and title:
                layout_size, layout_color = _read_layout_title_style(slide, idx)
                size_pt = layout_size or title_fallback
                # Clamp to fit placeholder height (leave ~30% for line spacing)
                if ph.height:
                    max_pt = int(Emu(ph.height).pt * 0.70)
                    if max_pt > 0:
                        size_pt = min(size_pt, max_pt)
                color = layout_color or _resolve_color("dk1", self.theme)

                ph.text = title
                tf = ph.text_frame
                tf.margin_left = 0
                tf.margin_top = 0
                tf.word_wrap = True
                for para in tf.paragraphs:
                    para.font.name = self.theme.fonts.major
                    para.font.bold = True
                    para.font.size = Pt(size_pt)
                    para.font.color.rgb = color
                filled.add(idx)

            elif is_subtitle_ph and subtitle:
                layout_size, layout_color = _read_layout_title_style(slide, idx)
                size_pt = layout_size or subtitle_fallback
                color = layout_color or _resolve_color("dk2", self.theme)

                ph.text = subtitle
                tf = ph.text_frame
                tf.margin_left = 0
                tf.margin_top = 0
                tf.word_wrap = True
                for para in tf.paragraphs:
                    para.font.name = self.theme.fonts.minor
                    para.font.size = Pt(size_pt)
                    para.font.color.rgb = color
                filled.add(idx)

        return filled

    def _render_content_slide(
        self, slide, layout: SlideLayout, content: OptimizedSlideContent
    ) -> None:
        """Render a content slide with shapes, charts, tables.

        Title goes into the template's PH0 (preserves template font/positioning).
        Key message goes into PH1 as subtitle.
        All decorative and content shapes are then added on top.
        """
        # Populate template placeholders first so the slide title uses
        # the template's own font size, color, and position.
        filled = self._populate_placeholders(slide, content.title, content.key_message)

        # Strip leftover body/content placeholders (e.g. 'Source:', 'Click to
        # edit text') whose inherited layout text would overlap our shapes.
        self._strip_unused_placeholders(slide, keep_idxs=filled)

        # Render each shape from the layout
        for spec in layout.shapes:
            if spec.shape_type == "chart":
                if content.chart_data:
                    self.visual_gen.add_chart(slide, content.chart_data, spec)
            elif spec.shape_type == "table":
                if content.table_data:
                    self.visual_gen.add_table(slide, content.table_data, spec)
            else:
                self.visual_gen.add_shape(slide, spec, self.theme)
